from os.path import dirname, join
from textx import metamodel_from_file
from textx.export import metamodel_export
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Pt
import matplotlib.pyplot as plt
import sys
import os
from email.message import EmailMessage
import ssl
import smtplib

#Program class holds all of the methods needed to interpret a valid .instruct file
class Program:
    def interpret(self, model):
        #bools for each method type
        isSlides = False
        isDocument = False
        isLesson = False
        isDistribution = False
        isGrading = False
        for c in model.statements:
            methodName = c.__class__.__name__ 
            
            #interprets a lesson method
            if methodName == "Lesson":
                lessonName = c.className
                isLesson = True
                for d in c.singleStatement:
                    try:
                        if "Lesson Type" == d.name:
                            if "Slides" == d.type:
                                isSlides = not isDocument
                                lecture = Presentation()
                            elif "Document" == d.type:
                                isDocument = not isSlides
                                lessonDocument = Document()
                            else:
                                raise Exception()
                    except:
                        print("Lesson Types currently supported are Slides and Document")
                        sys.exit(1)
                    if "Slide Amount" == d.name:
                        try:
                            slideAmount = int(d.type)
                        except:
                            print("Type Error: " + d.type + " is not an integer")
                            sys.exit(1)
                    elif "File Name" == d.name:
                        exportedFileName = d.type
            
            #interprets a Slides method
            elif methodName == "Slides":
                if isSlides:
                    title_slide_layout = lecture.slide_layouts[0]
                    topic_slide_layout = lecture.slide_layouts[1]
                    try:
                        if len(c.singleStatement) != slideAmount:
                            raise Exception()
                    except:
                        print("Number of Slides specified does not match number of slides in Content")
                        sys.exit(1)
                    for d in c.singleStatement:
                        if "Title" in d.name:
                            slide = lecture.slides.add_slide(title_slide_layout)
                            title = slide.shapes.title
                            title.text = d.type
                            for shape in slide.shapes:
                                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                                    slide.shapes._spTree.remove(shape._element)
                        else:
                            slide = lecture.slides.add_slide(topic_slide_layout)
                            shapes = slide.shapes
                            title_shape = shapes.title
                            title_shape.text = d.type
            
            #interprets a Document Content method
            elif methodName == "DocumentContent":
                if isDocument:
                    for d in c.singleStatement:
                        if "Title" in d.name:
                            lessonDocument.add_heading(d.type, 0)
                        elif "Paragraph" in d.name and "Style" in d.name:
                            listItems = d.type.list.listItems
                            paragraphRun = currentParagraph.runs[0]
                            if listItems[0] == "Bold":
                                paragraphRun.bold = True
                            if "pt" in listItems[1]:
                                fontSizeString = listItems[1].split("pt")[0]
                                try:
                                    fontSize = int(fontSizeString)
                                    paragraphRun.font.size = Pt(fontSize)
                                except:
                                    print("Error: font size needs to have an integer value associated with it")
                                    sys.exit(1)
                            if len(listItems) == 3:
                                if isinstance(listItems[2], str):
                                    try:
                                        fontType = listItems[2]
                                        paragraphRun.font.name = fontType
                                    except:
                                        print("Font Type is invalid or cannot be found")
                                        sys.exit(1)
                        elif "Paragraph" in d.name:
                            currentParagraph = lessonDocument.add_paragraph(d.type)
                        elif "Add" in d.name and "Empty Lines" in d.name:
                            emptyLineString = d.name.split("Add ")[1].split( "Empty Lines")[0]
                            try:
                                emptyLineAmount = int(emptyLineString)
                                i = 0
                                while i < emptyLineAmount:
                                    lessonDocument.add_paragraph()
                                    i += 1
                            except:
                                print("Empty Line Amout should be an integer value")
                                sys.exit(1)
            
            #interprets a Slide Content method
            elif methodName == "SlideContent":
                belowTitleSlideTitle = [Inches(2.9),Inches(3.75)]
                rightHalf = [Inches(4.67), Inches(2.5)]
                leftHalf = [Inches(0), Inches(2.5)]
                height = Inches(3) 
                slideWidth = lecture.slide_width
                for d in c.singleStatement:
                    slideContentName = d.name
                    if "Slide" and "Image" and "Location" in slideContentName:
                        if "Below Title Slide Title" in d.type:
                            isRightHalf = False
                            isLeftHalf = False
                            left = belowTitleSlideTitle[0]
                            top = belowTitleSlideTitle[1]
                        elif "Right Half" in d.type:
                            isRightHalf = True
                            isLeftHalf = False
                            left = rightHalf[0]
                            top = rightHalf[1]
                        elif "Left Half" in d.type:
                            isRightHalf = False
                            isLeftHalf = True
                            left = leftHalf[0]
                            top = leftHalf[1]
                    elif "Slide" and "Image" in slideContentName:
                        slideContentName = slideContentName.split("Slide ")[1].split( "Image")[0]
                        slideContentIndex = int(slideContentName)
                        currentSlide = lecture.slides[slideContentIndex - 1]
                        currentImage = currentSlide.shapes.add_picture(d.type, left, top, height=height)
                        if isRightHalf:
                            currentImage.left = slideWidth - currentImage.width
                    elif "Slide" and "Bullet Points" in slideContentName:
                        slideContentName = slideContentName.split("Slide ")[1].split( "Bullet Points")[0]
                        slideContentIndex = int(slideContentName)
                        currentSlide = lecture.slides[slideContentIndex - 1]
                        bulletPointList = d.type.list.listItems
                        
                        textBoxPlaceholder = currentSlide.shapes.placeholders[1]
                        textBoxPlaceholder.width = Inches(slideWidth.inches / 2)
                        textBoxPlaceholder.top = Inches(1)
                        textFrame = textBoxPlaceholder.text_frame
                        if isLeftHalf:
                            textBoxPlaceholder.left = slideWidth - textBoxPlaceholder.width
                            self.AddBulletPoints(bulletPointList,textFrame)
                        elif isRightHalf:
                            textBoxPlaceholder.left = Inches(0)
                            self.AddBulletPoints(bulletPointList,textFrame)
            
            #interprets a grading method
            elif methodName == "Grading":
                testName = c.className
                isGrading = True
                isLesson = not isGrading
                isDistribution = not isGrading
                isMC = False
                isTF = False
                for d in c.singleStatement:
                    if "Test Type" == d.name:
                        testType = d.type
                        if testType == "Multiple Choice":
                            isMC = True
                            isTF = False
                        elif testType == "True/False":
                            isMC = False
                            isTF = True
                    elif "Number of Questions" == d.name:
                        numQuestions = d.type
                    elif "Grade Report File" == d.name:
                        gradeReportFile = d.type
                    elif "Recipient Email" == d.name:
                        recipientEmail = d.type.email
                    elif "Grading Scale" == d.name:
                        gradingScaleName = d.type
            
            #interprets a Test method
            elif methodName == "Test":
                statements = c.singleStatement
                if isMC:
                    numCorrectVsNumTotal, correctToTotalRatio, incorrectAnswers = self.GradeTest(numQuestions,statements)
                elif isTF:
                    numCorrectVsNumTotal, correctToTotalRatio, incorrectAnswers = self.GradeTest(numQuestions,statements)
               
            #interprets a Grading Scale Method
            elif methodName == "GradingScale":
                if gradingScaleName == c.className:
                    gradingScale = {}
                    for d in c.singleStatement:
                        gradingScale.update({d.type: d.name})
            
            #interprets a Data Visualization method
            elif methodName == "DataVisualization":
                dataVizName = c.className
                isDistribution = True
                isGrading = not isDistribution
                isLesson = not isDistribution
                gradeAmounts = []
                for d in c.singleStatement:
                    match d.name:
                        case "GraphType":
                           graphType = d.type
                        case "Title":
                            graphTitle = d.type
                        case "X Axis Title":
                            xAxisTitle = d.type
                        case "Y Axis Title":
                            yAxisTitle = d.type
                        case "X Axis":
                           xAxisItems = d.type.list.listItems
                        case "Number of Total Students":
                            totalNumStudents = d.type
                        case "Number of A Grades":
                           numAGrades = d.type
                           gradeAmounts.append(numAGrades)
                        case "Number of B Grades":
                           numBGrades = d.type
                           gradeAmounts.append(numBGrades)
                        case "Number of C Grades":
                           numCGrades = d.type
                           gradeAmounts.append(numCGrades)
                        case "Number of D Grades":
                           numDGrades = d.type
                           gradeAmounts.append(numDGrades)
                        case "Number of F Grades":
                           numFGrades = d.type
                           gradeAmounts.append(numFGrades)
                        case "File Name":
                           graphFileType = d.type
                        case _:
                            continue
                if graphType == "Bar Graph":
                    isBarGraph = True
                    isPieChart = False
                    barGraphCategories = xAxisItems
                    barGraphValues = gradeAmounts
                    plt.bar(barGraphCategories, barGraphValues)
                    plt.title(graphTitle)
                    plt.xlabel(xAxisTitle)
                    plt.ylabel(yAxisTitle)
                elif graphType == "Pie Chart":
                    isPieChart = True
                    isBarGraph = False
                    try:
                        if sum(gradeAmounts) != totalNumStudents:
                            raise Exception()
                    except:
                        print("Error in method, " + methodName + " " + dataVizName +": Total Number of Students must match the sum of the grades in each category")
                        sys.exit(1)
                    pieChartLabels = ["A", "B", "C", "D", "F"]
                    pieChartSizes = self.PieSizes(gradeAmounts,totalNumStudents)
                    
            #interprets single line commands that aren't assignments statements or methods            
            elif methodName == "Command":
                commandString = c.name
                
                #interprets commands for a lesson template
                if(isLesson):
                    if "Create" and lessonName in commandString:
                        if isSlides:
                            lecture.save(exportedFileName)
                        elif isDocument:
                            lessonDocument.save(exportedFileName)
                    elif "Create" in commandString and lessonName not in commandString:
                        try:
                            incorrectLessonName = commandString.split(" Create")[1]
                            raise Exception()
                        except:
                            "Cannot create " + "\"" + incorrectLessonName + "\". It has not been created yet."
                            sys.exit(1)
                
                #interprets commands for a grading program
                elif(isGrading):
                    if "Grade" in commandString and testName in commandString:
                        with open(gradeReportFile,"w") as file:
                            file.write(testName + " Results: \n")
                            file.write(str(numCorrectVsNumTotal) + ", " + str(correctToTotalRatio * 100) + "% " + self.GradeMatch(gradingScale, correctToTotalRatio) + "\n")
                            if len(incorrectAnswers) > 0:
                                if len(incorrectAnswers) == 1:
                                    file.write("Answer to Question " + str(incorrectAnswers[0]) + " is incorrect.")
                                else:
                                    file.write("Answers to Questions ")
                                    for element in incorrectAnswers:
                                        file.write(str(element) + " ")
                                    file.write("are incorrect.")
                    elif "Grade" in commandString and testName not in commandString:
                        try:
                            incorrectGradeName = commandString.split("Grade ")[1]
                            raise Exception()
                        except:
                            print("Cannot grade " + "\"" + incorrectGradeName + "\". It has not been created yet.")
                            sys.exit(1)
                    elif "Send Email For" in commandString and testName in commandString:
                        with open(gradeReportFile, "r") as file:
                            emailBody = file.read()
                        self.SendEmail(recipientEmail,emailBody,testName)
                
                #interprets commands for a grade distribution program
                elif(isDistribution):
                    if "Create" and dataVizName in commandString:
                        if isBarGraph:
                            plt.savefig(graphFileType, format='jpg')
                            plt.close()
                        elif isPieChart:
                            plt.pie(pieChartSizes, labels=pieChartLabels, autopct='%1.1f%%', startangle=140)
                            plt.title(graphTitle)
                            plt.savefig(graphFileType, format='jpg')
                            plt.close()
                    elif "Create" in commandString and dataVizName not in commandString:
                        incorrectDataVizName = commandString.split("Create ")[1]
                        try:
                            raise Exception()
                        except:
                            print("Cannot create " + "\"" + incorrectDataVizName + "\". It has not been created yet.")
                            sys.exit(1)

    #helper method for checking if two lists were created in file .instruct file and if they are the same size
    #if so it will compare answers for output                  
    def GradeTest(self, numQuestions, statement):
        listOfListStatements = statement
        statementException = "Test needs two statements, one for Correct Answers, the other for Student Answers. Recieved: " + listOfListStatements[0].name + " and " + listOfListStatements[1].name
        if len(listOfListStatements) != 2:
            try:
                raise Exception()
            except:
                print("Only two statements are allowed in Test: Correct Answers and Student Answers. Recieved " + str(len(listOfListStatements)))
                sys.exit(1)
        else:
            if listOfListStatements[0].name == "Correct Answers":
                correctAnswersList = listOfListStatements[0].type.list.listItems
                if len(correctAnswersList) != numQuestions:
                    try:
                        raise Exception()
                    except:
                        print("Correct Answers must be same length as specified number of questions. Number of Questions = " + str(numQuestions) + ", Length of Correct Answers: " + str(len(correctAnswersList)))
                        sys.exit(1)
                if listOfListStatements[1].name == "Student Answers":
                    studentAnswersList = listOfListStatements[1].type.list.listItems
                    if len(studentAnswersList) != len(correctAnswersList):
                        try:
                            raise Exception()
                        except:
                            print("Length of Student Answers must be the same as Correct Answers. Correct Answers Length = " + str(len(correctAnswersList)) + " Student Answers Length = " + str(len(studentAnswersList)))
                            sys.exit(1)
                    correctAnswers,incorrectAnswers = self.CompareAnswers(correctAnswersList,studentAnswersList)
                    numCorrectVsNumTotal = str(correctAnswers) + "/" + str(len(correctAnswersList))
                    correctToTotalRatio = float(correctAnswers) / float(len(correctAnswersList))
                else:
                    print(statementException)
                    sys.exit(1)
            elif listOfListStatements[0].name == "Student Answers":
                studentAnswersList = listOfListStatements[0].type.list.listItems
                if listOfListStatements[1].name == "Correct Answers":
                    correctAnswersList = listOfListStatements[1].type.list.listItems
                    correctAnswers,incorrectAnswers = self.CompareAnswers(correctAnswersList,studentAnswersList)
                    numCorrectVsNumTotal = str(correctAnswers) + "/" + str(len(correctAnswersList))
                    correctToTotalRatio = float(correctAnswers) / float(len(correctAnswersList))
                else:
                    print(statementException)
                    sys.exit(1)
            else:
                print(statementException)
                sys.exit(1)
        return numCorrectVsNumTotal, correctToTotalRatio, incorrectAnswers
    
    #used by GradeTest(), it will compare each answer in each answer list to return the number of incorrect and correct answers
    def CompareAnswers(self, correctAnswers, studentAnswers):
        index = 0
        numCorrectAnswers= 0
        incorrectAnswers = []
        if len(correctAnswers) == len(studentAnswers):
            for answer in studentAnswers:
                if answer == correctAnswers[index]:
                    numCorrectAnswers += 1
                else:
                    incorrectAnswers.append(index)
                index += 1
        return numCorrectAnswers, incorrectAnswers
    
    #returns the letter grade based on the grading scale given in the .instruct file
    def GradeMatch(self, gradingScale, ratio):
        gradeThresholds = sorted(gradingScale.keys())
        gradeThresholdsTail = gradeThresholds[len(gradeThresholds) - 1]
        currentIndex = 0
        for x in gradeThresholds:
            if (ratio * 100) < x:
                return gradingScale[x]
            elif(ratio * 100) >= x and currentIndex == len(gradeThresholds) - 1 :
                return gradingScale[gradeThresholdsTail]
            currentIndex += 1
    
    #creates the labels for the pie chart based on the percentages of students who got a specific grade
    def PieSizes(self, amounts, totalNum):
        pieSizes = []
        for x in amounts:
            pieSizes.append(x / totalNum)
        return pieSizes
    
    #adds bullet points to a slide
    def AddBulletPoints(self,bulletPoints,textFrame):
        bulletPointIndex = 0
        while bulletPointIndex < len(bulletPoints):
            currentParagraph = textFrame.add_paragraph()
            currentParagraph.text = bulletPoints[bulletPointIndex]
            currentParagraph.level = 0
            bulletPointIndex += 1

    #emails a grade report to a recipient email specified in the .instruct file using the smtp, ssl, and emailMessage libraries. 
    #the email sender's address is instructifytest.@gmail.com
    def SendEmail(self, receiverAddress, emailBody, subjectLine):
        email_sender = "instructifytest@gmail.com"
        email_password = os.environ.get("EMAIL_PASSWORD")
        email_receiver = receiverAddress


        subject = subjectLine + " Grade"
        body = emailBody

        em = EmailMessage()
        em["From"] = email_sender
        em["To"] = email_receiver
        em["Subject"] = subject
        em.set_content(body)

        context = ssl.create_default_context()

        with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as smtp:
            smtp.login(email_sender, email_password)
            smtp.sendmail(email_sender,email_receiver, em.as_string())

#main method which uses textx to parse the .instruct file before it is interpreted semantically
def main(debug=False):

    this_folder = dirname(__file__)

    instructify_mm = metamodel_from_file(join(this_folder, 'instructify.tx'), debug=False)

    #example programs (can be replaced with custom named programs)
    #to run a particular program change the index inside programs[]. indexes are 0-(length of programs - 1)
    programs = ["lessonTemplate.instruct", "docTemplate.instruct", "gradeTest.instruct", "gradeDistribution.instruct"]
    instructify_model = instructify_mm.model_from_file(join(this_folder, programs[3]))

    program = Program()
    program.interpret(instructify_model)

if __name__ == "__main__":
    main()